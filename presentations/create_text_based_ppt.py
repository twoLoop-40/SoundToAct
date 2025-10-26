#!/usr/bin/env python3
"""
PowerPoint Generator - Text-Based Visuals

Creates presentation with text descriptions instead of images.
Includes detailed speaker notes for presenters.
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Color palette
COLORS = {
    'primary': RGBColor(37, 99, 235),
    'secondary': RGBColor(124, 58, 237),
    'accent': RGBColor(16, 185, 129),
    'warning': RGBColor(245, 158, 11),
    'text_primary': RGBColor(17, 24, 39),
    'text_secondary': RGBColor(107, 114, 128),
    'bg_light': RGBColor(240, 245, 255),
    'bg_accent': RGBColor(236, 253, 245),
}


def create_text_box_with_style(slide, left, top, width, height, text, font_size=18, bold=False, bg_color=None):
    """Create a styled text box"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = COLORS['text_primary']
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    return text_box


def create_visual_text_box(slide, left, top, width, height, text):
    """Create a visual description text box with border"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['bg_light']
    shape.line.color.rgb = COLORS['primary']
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text_primary']
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    return shape


def create_title_slide(prs, slide_data):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if slide_data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['subtitle']
        p.font.size = Pt(36)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER

    # Content
    content = slide_data.get('content', [])
    if content:
        content_top = Inches(4.5)
        for line in content:
            if line.strip():
                text_box = slide.shapes.add_textbox(Inches(1), content_top, Inches(8), Inches(0.4))
                tf = text_box.text_frame
                p = tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = COLORS['text_primary']
                p.alignment = PP_ALIGN.CENTER
                content_top += Inches(0.45)

    # Visuals as text boxes
    visuals = slide_data.get('visuals', [])
    if visuals:
        visual_top = Inches(5.5)
        for visual in visuals[:3]:  # Max 3 for title slide
            create_visual_text_box(slide, Inches(1.5), visual_top, Inches(7), Inches(0.5), visual)
            visual_top += Inches(0.6)

    # Speaker notes
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data['speaker_notes']

    return slide


def create_visual_slide(prs, slide_data):
    """Create slide with visual descriptions"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    current_top = Inches(1)

    # Subtitle
    if slide_data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(Inches(0.5), current_top, Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['subtitle']
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER
        current_top += Inches(0.6)

    # Content
    content = slide_data.get('content', [])
    if content and any(c.strip() for c in content):
        for line in content:
            if line.strip():
                text_box = slide.shapes.add_textbox(Inches(1), current_top, Inches(8), Inches(0.4))
                tf = text_box.text_frame
                p = tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = COLORS['text_primary']
                p.alignment = PP_ALIGN.CENTER
                current_top += Inches(0.45)
        current_top += Inches(0.3)

    # Visuals as text boxes
    visuals = slide_data.get('visuals', [])
    if visuals:
        num_visuals = len(visuals)

        # Calculate layout based on number of visuals
        if num_visuals <= 3:
            # Vertical stack
            box_width = Inches(8)
            box_left = Inches(1)
            box_height = Inches(0.6)

            for visual in visuals:
                if visual.strip():  # Skip empty lines
                    create_visual_text_box(slide, box_left, current_top, box_width, box_height, visual)
                    current_top += box_height + Inches(0.1)

        else:
            # Grid layout (2 columns)
            box_width = Inches(4.3)
            box_height = Inches(0.6)

            for i, visual in enumerate(visuals):
                if visual.strip():
                    row = i // 2
                    col = i % 2
                    box_left = Inches(0.5 + col * 4.8)
                    box_top = current_top + row * (box_height + Inches(0.1))

                    create_visual_text_box(slide, box_left, box_top, box_width, box_height, visual)

    # Speaker notes
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data['speaker_notes']

    return slide


def create_presentation_from_json(json_path, output_path):
    """Create PowerPoint presentation from JSON"""

    # Load JSON data
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    metadata = data['metadata']
    slides_data = data['slides']

    print(f"Creating presentation: {metadata['title']}")
    print(f"Total slides: {metadata['total_slides']}")
    print()

    # Create each slide
    for slide_data in slides_data:
        slide_num = slide_data['number']
        layout = slide_data['layout']

        print(f"  Creating slide {slide_num}: {slide_data['title']} ({layout})")

        if layout == 'TitleSlide':
            slide = create_title_slide(prs, slide_data)
        else:
            slide = create_visual_slide(prs, slide_data)

    # Save presentation
    prs.save(output_path)
    print(f"\n✓ Presentation saved: {output_path}")

    return output_path


def main():
    """Main function"""
    base_dir = Path(__file__).parent
    json_path = base_dir / 'output' / 'presentation_from_idris.json'
    output_path = base_dir / 'output' / 'SoundToAct_TextBased_Presentation.pptx'

    if not json_path.exists():
        print(f"Error: JSON file not found: {json_path}")
        return

    print("="*70)
    print("PowerPoint Generator - Text-Based Visuals")
    print("="*70)
    print()

    create_presentation_from_json(json_path, output_path)

    print()
    print("="*70)
    print("DONE!")
    print("="*70)
    print(f"\nYou can now open: {output_path}")
    print()
    print("📝 발표자 가이드:")
    print("   - 각 슬라이드의 발표자 노트를 꼭 확인하세요")
    print("   - PowerPoint에서 '노트' 보기로 전환하면 상세 가이드가 표시됩니다")
    print("   - 기술적 내용이 포함되어 있으니 충분히 숙지하세요")


if __name__ == '__main__':
    main()
