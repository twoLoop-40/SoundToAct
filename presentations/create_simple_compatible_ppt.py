#!/usr/bin/env python3
"""
PowerPoint Generator - Maximum Keynote Compatibility
Uses built-in slide layouts instead of custom shapes
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation_simple(json_path, output_path):
    """Create PowerPoint with maximum compatibility"""

    # Load JSON data
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Create presentation with default template (better Keynote compatibility)
    prs = Presentation()

    slides_data = data['slides']

    print(f"Creating presentation with {len(slides_data)} slides...")
    print()

    for slide_data in slides_data:
        slide_num = slide_data['number']
        print(f"  Creating slide {slide_num}: {slide_data['title']}")

        # Use built-in layout: Title and Content (layout 1)
        slide_layout = prs.slide_layouts[5]  # Blank layout for more control
        slide = prs.slides.add_slide(slide_layout)

        # Clear any existing shapes
        for shape in slide.shapes:
            if hasattr(shape, "element"):
                shape.element.getparent().remove(shape.element)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.3),
            Inches(9),
            Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = slide_data['title']
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(37, 99, 235)

        # Add subtitle if exists
        current_y = Inches(1.2)
        if slide_data.get('subtitle'):
            sub_box = slide.shapes.add_textbox(
                Inches(0.5),
                current_y,
                Inches(9),
                Inches(0.5)
            )
            sub_frame = sub_box.text_frame
            sub_frame.word_wrap = True
            p = sub_frame.paragraphs[0]
            p.text = slide_data['subtitle']
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(107, 114, 128)
            current_y += Inches(0.6)

        # Add content lines
        content = slide_data.get('content', [])
        for line in content:
            if line and line.strip():
                line_box = slide.shapes.add_textbox(
                    Inches(1),
                    current_y,
                    Inches(8),
                    Inches(0.4)
                )
                line_frame = line_box.text_frame
                line_frame.word_wrap = True
                p = line_frame.paragraphs[0]
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(17, 24, 39)
                current_y += Inches(0.45)

        # Add visuals as simple text boxes
        current_y += Inches(0.2)
        visuals = slide_data.get('visuals', [])

        for visual in visuals:
            if visual and visual.strip():
                # Simple text box - no shapes, no borders
                vis_box = slide.shapes.add_textbox(
                    Inches(1.5),
                    current_y,
                    Inches(7),
                    Inches(0.35)
                )
                vis_frame = vis_box.text_frame
                vis_frame.word_wrap = True
                p = vis_frame.paragraphs[0]
                p.text = visual
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(37, 99, 235)
                current_y += Inches(0.38)

        # Add speaker notes
        if slide_data.get('speaker_notes'):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data['speaker_notes']

    # Save
    prs.save(output_path)
    print(f"\n✓ Presentation saved: {output_path}")
    return output_path


def main():
    """Main function"""
    base_dir = Path(__file__).parent
    json_path = base_dir / 'output' / 'presentation_from_idris.json'
    output_path = base_dir / 'output' / 'SoundToAct_Simple.pptx'

    if not json_path.exists():
        print(f"Error: JSON file not found: {json_path}")
        return

    print("="*70)
    print("PowerPoint Generator - Simple & Compatible")
    print("="*70)
    print()

    create_presentation_simple(json_path, output_path)

    print()
    print("="*70)
    print("DONE!")
    print("="*70)
    print(f"\nFile: {output_path}")
    print("\nThis version uses:")
    print("  - Default PowerPoint template")
    print("  - Simple text boxes only")
    print("  - No custom shapes or borders")
    print("  - Maximum Keynote compatibility")


if __name__ == '__main__':
    main()
