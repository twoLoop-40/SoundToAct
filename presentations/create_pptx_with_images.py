#!/usr/bin/env python3
"""
PowerPoint Generator with Real Images

Creates presentation with actual images inserted into slides.
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Color palette
COLORS = {
    'primary': RGBColor(37, 99, 235),
    'secondary': RGBColor(124, 58, 237),
    'accent': RGBColor(16, 185, 129),
    'warning': RGBColor(245, 158, 11),
    'text_primary': RGBColor(17, 24, 39),
    'text_secondary': RGBColor(107, 114, 128),
}


def create_title_slide(prs, slide_data, images_dir):
    """Create title slide with images"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Map visual descriptions to image files
    image_map = {
        '음성 웨이브폼 애니메이션': 'waveform_animation.png',
        '마이크 아이콘 (큼직하게)': 'microphone_large.png',
        '만화 스타일 일러스트: 침대에서 일어나는 학생': 'student_waking_up.png',
        '복잡한 과정 플로우: 폰 찾기 → 잠금 해제 → 연락처 앱 → 검색 → 터치': 'complex_process_flow.png',
        '시계 아이콘: \'2분 소요\'': 'clock_two_minutes.png',
        '큰 물음표 아이콘': 'question_mark_large.png',
        '말풍선 안에 \'엄마\'': 'speech_bubble_mom.png',
        '빛나는 효과 (반짝이는 전구)': 'light_bulb_sparkle.png',
        '프로젝트 로고 (크게)': 'soundtoact_logo_large.png',
        '10초 데모 영상: \'엄마\' → 전화 걸림': 'demo_video_placeholder.png',
        'Before/After 비교 이미지': 'before_after_comparison.png',
        '1단계: 듣기 - 마이크 아이콘 + 음성 웨이브': 'step1_listen.png',
        '2단계: 이해하기 - AI 뇌 + 키워드 매칭': 'step2_understand.png',
        '3단계: 실행하기 - 액션 아이콘 (전화, 음악, 조명)': 'step3_act.png',
        '화살표로 연결된 3단계 플로우': 'three_step_flow_arrows.png',
        '실제 사용 데모 영상 (30초)': 'live_demo_video.png',
        '데모 스크린샷 (백업)': 'demo_screenshot_backup.png',
        'Before: 복잡한 과정 (2분)': 'before_2_minutes.png',
        'After: 말 한마디 (2초)': 'after_2_seconds.png',
        '숫자 강조: 60배 빨라짐': 'sixty_times_faster.png',
        '하루 30분 절약': 'thirty_minutes_saved.png',
        '시나리오 1: 어르신 - 큰 글씨 필요없이': 'elderly_scenario.png',
        '시나리오 2: 바쁜 직장인 - 운전 중에도': 'worker_scenario.png',
        '시나리오 3: 장애인 - 손 사용 불편해도': 'disability_scenario.png',
        '모두를 위한 기술': 'inclusive_technology.png',
        '지구 아이콘 + 연결된 사람들': 'connected_world.png',
        '밝은 미래 이미지': 'bright_future.png',
        '확장 가능성: 스마트홈, 자동차, 가전제품...': 'expansion_vision.png',
        'QR 코드 (GitHub)': 'qr_code_github.png',
        'SoundToAct 로고': 'soundtoact_logo_final.png',
    }

    # Title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    current_top = Inches(3.5)
    if slide_data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(left, current_top, width, Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['subtitle']
        p.font.size = Pt(36)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER
        current_top = Inches(4.5)

    # Process visuals
    visuals = slide_data.get('visuals', [])
    num_visuals = len(visuals)

    if num_visuals > 0:
        if num_visuals == 1:
            # Single visual centered
            visual = visuals[0]
            img_file = None
            if visual in image_map:
                img_path = images_dir / image_map[visual]
                if img_path.exists():
                    img_file = img_path

            if img_file:
                slide.shapes.add_picture(
                    str(img_file),
                    Inches(3), current_top,
                    width=Inches(4)
                )
                img_bottom = current_top + Inches(3)
                # Caption
                caption_box = slide.shapes.add_textbox(Inches(1), img_bottom + Inches(0.1), width, Inches(0.4))
                tf = caption_box.text_frame
                p = tf.paragraphs[0]
                p.text = visual
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['text_secondary']
                p.alignment = PP_ALIGN.CENTER

        elif num_visuals == 2:
            # Two visuals side by side
            for i, visual in enumerate(visuals):
                img_file = None
                if visual in image_map:
                    img_path = images_dir / image_map[visual]
                    if img_path.exists():
                        img_file = img_path

                box_left = Inches(1 + i * 4.5)
                box_width = Inches(4)

                if img_file:
                    slide.shapes.add_picture(
                        str(img_file),
                        box_left, current_top,
                        width=box_width
                    )
                    img_bottom = current_top + Inches(2.5)
                    # Caption
                    caption_box = slide.shapes.add_textbox(box_left, img_bottom + Inches(0.1), box_width, Inches(0.4))
                    tf = caption_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = visual
                    p.font.size = Pt(12)
                    p.font.color.rgb = COLORS['text_secondary']
                    p.alignment = PP_ALIGN.CENTER

        else:
            # Three visuals in a row
            for i, visual in enumerate(visuals[:3]):
                img_file = None
                if visual in image_map:
                    img_path = images_dir / image_map[visual]
                    if img_path.exists():
                        img_file = img_path

                box_left = Inches(0.75 + i * 3.25)
                box_width = Inches(3)

                if img_file:
                    slide.shapes.add_picture(
                        str(img_file),
                        box_left, current_top,
                        width=box_width
                    )
                    img_bottom = current_top + Inches(2)
                    # Caption
                    caption_box = slide.shapes.add_textbox(box_left, img_bottom + Inches(0.05), box_width, Inches(0.4))
                    tf = caption_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = visual
                    p.font.size = Pt(11)
                    p.font.color.rgb = COLORS['text_secondary']
                    p.alignment = PP_ALIGN.CENTER

    # Content at bottom
    content = slide_data.get('content', [])
    if content:
        text_box = slide.shapes.add_textbox(left, Inches(6.5), width, Inches(0.5))
        tf = text_box.text_frame
        for i, line in enumerate(content):
            if line.strip():
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = COLORS['text_primary']
                p.alignment = PP_ALIGN.CENTER

    return slide


def create_visual_slide(prs, slide_data, images_dir):
    """Create slide with large images"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle if exists
    current_top = Inches(1.1)
    if slide_data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(left, current_top, width, Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['subtitle']
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER
        current_top += Inches(0.6)

    # Map visual descriptions to image files (same as create_title_slide)
    image_map = {
        '음성 웨이브폼 애니메이션': 'waveform_animation.png',
        '마이크 아이콘 (큼직하게)': 'microphone_large.png',
        '만화 스타일 일러스트: 침대에서 일어나는 학생': 'student_waking_up.png',
        '복잡한 과정 플로우: 폰 찾기 → 잠금 해제 → 연락처 앱 → 검색 → 터치': 'complex_process_flow.png',
        '시계 아이콘: \'2분 소요\'': 'clock_two_minutes.png',
        '큰 물음표 아이콘': 'question_mark_large.png',
        '말풍선 안에 \'엄마\'': 'speech_bubble_mom.png',
        '빛나는 효과 (반짝이는 전구)': 'light_bulb_sparkle.png',
        '프로젝트 로고 (크게)': 'soundtoact_logo_large.png',
        '10초 데모 영상: \'엄마\' → 전화 걸림': 'demo_video_placeholder.png',
        'Before/After 비교 이미지': 'before_after_comparison.png',
        '1단계: 듣기 - 마이크 아이콘 + 음성 웨이브': 'step1_listen.png',
        '2단계: 이해하기 - AI 뇌 + 키워드 매칭': 'step2_understand.png',
        '3단계: 실행하기 - 액션 아이콘 (전화, 음악, 조명)': 'step3_act.png',
        '화살표로 연결된 3단계 플로우': 'three_step_flow_arrows.png',
        '실제 사용 데모 영상 (30초)': 'live_demo_video.png',
        '데모 스크린샷 (백업)': 'demo_screenshot_backup.png',
        'Before: 복잡한 과정 (2분)': 'before_2_minutes.png',
        'After: 말 한마디 (2초)': 'after_2_seconds.png',
        '숫자 강조: 60배 빨라짐': 'sixty_times_faster.png',
        '하루 30분 절약': 'thirty_minutes_saved.png',
        '시나리오 1: 어르신 - 큰 글씨 필요없이': 'elderly_scenario.png',
        '시나리오 2: 바쁜 직장인 - 운전 중에도': 'worker_scenario.png',
        '시나리오 3: 장애인 - 손 사용 불편해도': 'disability_scenario.png',
        '모두를 위한 기술': 'inclusive_technology.png',
        '지구 아이콘 + 연결된 사람들': 'connected_world.png',
        '밝은 미래 이미지': 'bright_future.png',
        '확장 가능성: 스마트홈, 자동차, 가전제품...': 'expansion_vision.png',
        'QR 코드 (GitHub)': 'qr_code_github.png',
        'SoundToAct 로고': 'soundtoact_logo_final.png',
    }

    visuals = slide_data.get('visuals', [])

    # Process all visuals
    num_visuals = len(visuals)

    if num_visuals == 0:
        return slide

    # Determine layout based on number of visuals
    if num_visuals == 1:
        # Single large visual (centered)
        visual = visuals[0]
        img_file = None
        if visual in image_map:
            img_path = images_dir / image_map[visual]
            if img_path.exists():
                img_file = img_path

        if img_file:
            # Display image
            img_width = Inches(7)
            img_left = Inches(1.5)
            img_top = current_top

            slide.shapes.add_picture(
                str(img_file),
                img_left, img_top,
                width=img_width
            )
            current_top += Inches(4.5) + Inches(0.1)

            # Add caption below image
            caption_box = slide.shapes.add_textbox(
                img_left, current_top, img_width, Inches(0.5)
            )
            tf = caption_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = visual
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['text_secondary']
            p.alignment = PP_ALIGN.CENTER
        else:
            # Placeholder box
            box_width = Inches(7)
            box_left = Inches(1.5)
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                box_left, current_top, box_width, Inches(3.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 245, 255)
            shape.line.color.rgb = COLORS['primary']
            shape.line.width = Pt(3)

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = visual
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text_secondary']
            p.alignment = PP_ALIGN.CENTER

    elif num_visuals == 2:
        # Two visuals side by side
        for i, visual in enumerate(visuals):
            img_file = None
            if visual in image_map:
                img_path = images_dir / image_map[visual]
                if img_path.exists():
                    img_file = img_path

            box_left = Inches(0.5 + i * 5)
            box_width = Inches(4.5)

            if img_file:
                # Display image
                slide.shapes.add_picture(
                    str(img_file),
                    box_left, current_top,
                    width=box_width
                )
                img_bottom = current_top + Inches(3)

                # Add caption below image
                caption_box = slide.shapes.add_textbox(
                    box_left, img_bottom + Inches(0.1), box_width, Inches(0.5)
                )
                tf = caption_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = visual
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['text_secondary']
                p.alignment = PP_ALIGN.CENTER
            else:
                # Placeholder box
                shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    box_left, current_top, box_width, Inches(3.5)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 245, 255)
                shape.line.color.rgb = COLORS['primary']
                shape.line.width = Pt(3)

                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = visual
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS['text_secondary']
                p.alignment = PP_ALIGN.CENTER

    else:
        # Grid layout for 3+ visuals (2x2 grid)
        for i, visual in enumerate(visuals[:4]):
            row = i // 2
            col = i % 2
            box_left = Inches(0.5 + col * 4.75)
            box_top = current_top + row * Inches(2.5)
            box_width = Inches(4.5)

            img_file = None
            if visual in image_map:
                img_path = images_dir / image_map[visual]
                if img_path.exists():
                    img_file = img_path

            if img_file:
                # Display image (smaller for grid)
                slide.shapes.add_picture(
                    str(img_file),
                    box_left, box_top,
                    width=box_width,
                    height=Inches(1.8)
                )
                img_bottom = box_top + Inches(1.8)

                # Add caption below image
                caption_box = slide.shapes.add_textbox(
                    box_left, img_bottom + Inches(0.05), box_width, Inches(0.6)
                )
                tf = caption_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = visual
                p.font.size = Pt(12)
                p.font.color.rgb = COLORS['text_secondary']
                p.alignment = PP_ALIGN.CENTER
            else:
                # Placeholder box
                shape = slide.shapes.add_shape(
                    1,
                    box_left, box_top, box_width, Inches(2.2)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 245, 255)
                shape.line.color.rgb = COLORS['primary']
                shape.line.width = Pt(2)

                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = visual
                p.font.size = Pt(12)
                p.font.color.rgb = COLORS['text_secondary']
                p.alignment = PP_ALIGN.CENTER

    # Content text at bottom
    content = slide_data.get('content', [])
    if content and any(c.strip() for c in content):
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.8)
        )
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(content):
            if line.strip():
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = COLORS['text_primary']
                p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation_from_json(json_path, images_dir, output_path):
    """Create PowerPoint presentation with images from JSON"""

    # Load JSON data
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    metadata = data['metadata']
    slides_data = data['slides']

    print(f"Creating presentation: {metadata['title']}")
    print(f"Total slides: {metadata['total_slides']}")
    print(f"Images directory: {images_dir}")
    print()

    # Create each slide
    for slide_data in slides_data:
        slide_num = slide_data['number']
        layout = slide_data['layout']

        print(f"  Creating slide {slide_num}: {slide_data['title']} ({layout})")

        # Choose slide creation function
        if layout == 'TitleSlide':
            slide = create_title_slide(prs, slide_data, images_dir)
        else:
            slide = create_visual_slide(prs, slide_data, images_dir)

    # Save presentation
    prs.save(output_path)
    print(f"\n✓ Presentation saved: {output_path}")

    return output_path


def main():
    """Main function"""
    base_dir = Path(__file__).parent
    json_path = base_dir / 'output' / 'presentation_from_idris.json'
    images_dir = base_dir / 'images'
    output_path = base_dir / 'output' / 'SoundToAct_Presentation_WithImages.pptx'

    if not json_path.exists():
        print(f"Error: JSON file not found: {json_path}")
        return

    if not images_dir.exists():
        print(f"Error: Images directory not found: {images_dir}")
        print("Please run generate_images.py first")
        return

    print("="*60)
    print("PowerPoint Generator with Images")
    print("="*60)
    print()

    create_presentation_from_json(json_path, images_dir, output_path)

    print()
    print("="*60)
    print("DONE!")
    print("="*60)
    print(f"\nYou can now open: {output_path}")


if __name__ == '__main__':
    main()
