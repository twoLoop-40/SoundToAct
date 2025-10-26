#!/usr/bin/env python3
"""
Test simple PowerPoint creation to verify library works
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_simple_test():
    """Create a very simple test presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add one simple slide
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Test Presentation"
    p.font.size = Pt(44)
    p.font.bold = True

    # Add text
    text_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "This is a simple test to verify PowerPoint generation works."
    p.font.size = Pt(24)

    # Save
    output_path = Path(__file__).parent / 'output' / 'TEST_Simple.pptx'
    prs.save(output_path)
    print(f"✓ Test file created: {output_path}")
    return output_path

if __name__ == '__main__':
    create_simple_test()
