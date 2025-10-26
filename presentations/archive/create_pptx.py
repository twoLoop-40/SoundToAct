#!/usr/bin/env python3
"""
PowerPoint Generator for SoundToAct Presentation

This script creates an actual PowerPoint file from the presentation specification.
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Color palette
COLORS = {
    'primary': RGBColor(37, 99, 235),      # #2563EB Blue
    'secondary': RGBColor(124, 58, 237),   # #7C3AED Purple
    'accent': RGBColor(16, 185, 129),      # #10B981 Green
    'warning': RGBColor(245, 158, 11),     # #F59E0B Orange
    'text_primary': RGBColor(17, 24, 39),  # #111827 Dark
    'text_secondary': RGBColor(107, 114, 128),  # #6B7280 Gray
}


def create_title_slide(prs, slide_data):
    """Create title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Set title
    title.text = "SoundToAct"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Set subtitle
    if slide_data.get('subtitle'):
        subtitle.text = slide_data['subtitle']
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLORS['text_secondary']

    # Add additional text from code blocks
    if slide_data.get('code_blocks'):
        tf = subtitle.text_frame
        for code_block in slide_data['code_blocks']:
            p = tf.add_paragraph()
            p.text = code_block.get('code', '')
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text_secondary']

    return slide


def create_content_slide(prs, slide_data):
    """Create content slide with title and VISUAL-FOCUSED design"""
    slide_layout = prs.slide_layouts[6]  # Blank layout for maximum flexibility
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle if exists
    if slide_data.get('subtitle'):
        sub_top = Inches(1.2)
        sub_box = slide.shapes.add_textbox(left, sub_top, width, Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['subtitle']
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER

    # VISUALS - Create large placeholder boxes
    visuals = slide_data.get('visuals', [])
    if visuals:
        visual_top = Inches(2.5)
        visual_height = Inches(4)

        if len(visuals) <= 2:
            # Large single or side-by-side visuals
            box_width = Inches(8) if len(visuals) == 1 else Inches(4)
            for i, visual in enumerate(visuals[:2]):
                box_left = Inches(1) if len(visuals) == 1 else Inches(0.5 + i * 4.5)
                shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    box_left, visual_top, box_width, visual_height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 240, 250)
                shape.line.color.rgb = COLORS['primary']
                shape.line.width = Pt(2)

                # Add visual description text
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"[{visual}]"
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = COLORS['text_secondary']
                p.alignment = PP_ALIGN.CENTER
        else:
            # Multiple visuals - grid layout
            visual_height = Inches(2.5)
            for i, visual in enumerate(visuals[:4]):
                row = i // 2
                col = i % 2
                box_left = Inches(0.5 + col * 4.75)
                box_top = Inches(2.2 + row * 2.7)

                shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    box_left, box_top, Inches(4.5), visual_height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 240, 250)
                shape.line.color.rgb = COLORS['primary']
                shape.line.width = Pt(2)

                # Add visual description
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"[{visual}]"
                p.font.size = Pt(14)
                p.font.italic = True
                p.font.color.rgb = COLORS['text_secondary']
                p.alignment = PP_ALIGN.CENTER

    # Minimal content text at the bottom if any
    content = slide_data.get('content', [])
    if content and len([c for c in content if c.strip()]) > 0:
        text_left = Inches(1)
        text_top = Inches(6.5)
        text_width = Inches(8)
        text_height = Inches(0.8)

        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(content):
            if line.strip():
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = COLORS['text_primary']
                p.alignment = PP_ALIGN.CENTER

    return slide


def create_two_column_slide(prs, slide_data):
    """Create two-column slide"""
    slide_layout = prs.slide_layouts[3]  # Two content layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = slide_data['title']
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Split content into two columns
    if slide_data.get('content'):
        mid = len(slide_data['content']) // 2

        # Left column
        left_content = slide.placeholders[1]
        tf = left_content.text_frame
        tf.clear()
        for i, line in enumerate(slide_data['content'][:mid]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

        # Right column
        right_content = slide.placeholders[2]
        tf = right_content.text_frame
        tf.clear()
        for i, line in enumerate(slide_data['content'][mid:]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

    return slide


def create_blank_slide(prs, slide_data):
    """Create blank slide for demos"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title in the center
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle if exists
    if slide_data.get('subtitle'):
        p = tf.add_paragraph()
        p.text = slide_data['subtitle']
        p.font.size = Pt(32)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_footer(slide, slide_number, total_slides):
    """Add footer with slide number"""
    left = Inches(8.5)
    top = Inches(7)
    width = Inches(1)
    height = Inches(0.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    p = tf.paragraphs[0]
    p.text = f"{slide_number}/{total_slides}"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_secondary']
    p.alignment = PP_ALIGN.RIGHT


def create_presentation_from_json(json_path, output_path):
    """Create PowerPoint presentation from JSON specification"""

    # Load JSON data
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    metadata = data['metadata']
    slides_data = data['slides']
    total_slides = metadata['total_slides']

    print(f"Creating presentation: {metadata['title']}")
    print(f"Total slides: {total_slides}")

    # Create each slide
    for slide_data in slides_data:
        slide_num = slide_data['number']
        layout = slide_data['layout']

        print(f"  Creating slide {slide_num}: {slide_data['title']} ({layout})")

        # Choose appropriate slide creation function based on layout
        if slide_num == 1 or layout == 'TitleSlide':
            slide = create_title_slide(prs, slide_data)
        elif layout == 'TwoColumn':
            slide = create_two_column_slide(prs, slide_data)
        elif layout == 'FullScreenDemo':
            slide = create_blank_slide(prs, slide_data)
        else:  # SingleColumn, ThreeColumn, FourQuadrant
            slide = create_content_slide(prs, slide_data)

        # Add footer (except for title and demo slides)
        if slide_num != 1 and layout != 'FullScreenDemo':
            add_footer(slide, slide_num, total_slides)

    # Save presentation
    prs.save(output_path)
    print(f"\n✓ Presentation saved: {output_path}")

    return output_path


def main():
    """Main function"""
    base_dir = Path(__file__).parent
    # Use the Idris-based JSON
    json_path = base_dir / 'output' / 'presentation_from_idris.json'
    output_path = base_dir / 'output' / 'SoundToAct_Presentation.pptx'

    if not json_path.exists():
        print(f"Error: JSON file not found: {json_path}")
        print("Please run slides_from_idris.py first to generate the JSON file.")
        return

    print("="*60)
    print("PowerPoint Generator for SoundToAct")
    print("="*60)
    print()

    create_presentation_from_json(json_path, output_path)

    print()
    print("="*60)
    print("DONE!")
    print("="*60)
    print(f"\nYou can now open: {output_path}")


if __name__ == '__main__':
    main()
