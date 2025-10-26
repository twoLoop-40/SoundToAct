#!/usr/bin/env python3
"""
PowerPoint Validator
Tests if the generated PPT matches the Idris specifications
"""

import json
from pathlib import Path
from pptx import Presentation


def validate_presentation(pptx_path, json_path, images_dir):
    """Validate PowerPoint against specifications"""

    print("="*70)
    print("PowerPoint Validation Against Specifications")
    print("="*70)
    print()

    # Load JSON spec
    with open(json_path, 'r', encoding='utf-8') as f:
        spec = json.load(f)

    # Load PowerPoint
    prs = Presentation(pptx_path)

    metadata = spec['metadata']
    slides_spec = spec['slides']

    print(f"Presentation: {metadata['title']}")
    print(f"Expected slides: {metadata['total_slides']}")
    print(f"Actual slides: {len(prs.slides)}")
    print()

    # Validation results
    issues = []
    checks_passed = 0
    checks_total = 0

    # Check 1: Slide count
    checks_total += 1
    if len(prs.slides) == metadata['total_slides']:
        print("✓ Slide count matches specification")
        checks_passed += 1
    else:
        issue = f"✗ Slide count mismatch: expected {metadata['total_slides']}, got {len(prs.slides)}"
        print(issue)
        issues.append(issue)

    print()
    print("-"*70)
    print("Validating individual slides:")
    print("-"*70)
    print()

    # Check each slide
    for i, slide_spec in enumerate(slides_spec):
        slide_num = slide_spec['number']
        expected_title = slide_spec['title']
        expected_layout = slide_spec['layout']
        expected_visuals = slide_spec.get('visuals', [])

        print(f"Slide {slide_num}: {expected_title}")

        if i >= len(prs.slides):
            issue = f"  ✗ Slide {slide_num} missing from presentation"
            print(issue)
            issues.append(issue)
            continue

        slide = prs.slides[i]

        # Check 2: Slide has content
        checks_total += 1
        if len(slide.shapes) > 0:
            print(f"  ✓ Has {len(slide.shapes)} shapes")
            checks_passed += 1
        else:
            issue = f"  ✗ Slide {slide_num} has no shapes"
            print(issue)
            issues.append(issue)

        # Check 3: Expected images exist
        image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)  # 13 = Picture
        if len(expected_visuals) > 0:
            checks_total += 1
            if image_count > 0:
                print(f"  ✓ Has {image_count} images (expected visuals: {len(expected_visuals)})")
                checks_passed += 1
            else:
                issue = f"  ✗ No images found (expected {len(expected_visuals)} visuals)"
                print(issue)
                issues.append(issue)

        # Check 4: Text content exists
        text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text')]
        checks_total += 1
        if text_shapes:
            all_text = ' '.join(shape.text for shape in text_shapes)
            if expected_title in all_text or any(expected_title.split()[0] in all_text for word in expected_title.split()):
                print(f"  ✓ Title text found in slide")
                checks_passed += 1
            else:
                issue = f"  ✗ Title text not found in slide"
                print(issue)
                issues.append(issue)
        else:
            issue = f"  ✗ No text shapes found"
            print(issue)
            issues.append(issue)

        print()

    # Check 5: All required images exist
    print("-"*70)
    print("Checking image files:")
    print("-"*70)
    print()

    required_images = [
        'wave_animation.png',
        'mic_icon_large.png',
        'morning_illustration.png',
        'complex_process.png',
        'question_mark.png',
        'speech_bubble_mom.png',
        'light_bulb.png',
        'logo.png',
        'demo_placeholder.png',
        'three_steps.png',
        'live_demo.png',
        'before_after.png',
        'inclusive_tech.png',
        'world_connections.png',
        'bright_future.png',
        'qr_code.png',
        'final_logo.png',
    ]

    for img_name in required_images:
        checks_total += 1
        img_path = images_dir / img_name
        if img_path.exists():
            size_kb = img_path.stat().st_size / 1024
            print(f"  ✓ {img_name} ({size_kb:.1f} KB)")
            checks_passed += 1
        else:
            issue = f"  ✗ Missing: {img_name}"
            print(issue)
            issues.append(issue)

    print()
    print("="*70)
    print("VALIDATION SUMMARY")
    print("="*70)
    print()
    print(f"Checks passed: {checks_passed}/{checks_total}")
    print(f"Success rate: {checks_passed/checks_total*100:.1f}%")
    print()

    if issues:
        print(f"Issues found: {len(issues)}")
        print()
        for issue in issues:
            print(issue)
        print()
        return False
    else:
        print("✓ All validations passed!")
        print()
        return True


def main():
    """Main validation function"""
    base_dir = Path(__file__).parent
    pptx_path = base_dir / 'output' / 'SoundToAct_Presentation_WithImages.pptx'
    json_path = base_dir / 'output' / 'presentation_from_idris.json'
    images_dir = base_dir / 'images'

    if not pptx_path.exists():
        print(f"Error: PowerPoint file not found: {pptx_path}")
        return False

    if not json_path.exists():
        print(f"Error: JSON spec not found: {json_path}")
        return False

    if not images_dir.exists():
        print(f"Error: Images directory not found: {images_dir}")
        return False

    success = validate_presentation(pptx_path, json_path, images_dir)

    if success:
        print("="*70)
        print("✓ VALIDATION SUCCESSFUL")
        print("="*70)
        print()
        print("The PowerPoint presentation matches the Idris2 specifications.")
    else:
        print("="*70)
        print("⚠ VALIDATION FAILED")
        print("="*70)
        print()
        print("The PowerPoint presentation has some issues.")
        print("Please review the issues above and fix them.")

    return success


if __name__ == '__main__':
    import sys
    success = main()
    sys.exit(0 if success else 1)
