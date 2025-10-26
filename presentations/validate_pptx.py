#!/usr/bin/env python3
"""
Validate PowerPoint file integrity
"""

from pathlib import Path
from pptx import Presentation
import json

def validate_pptx(pptx_path):
    """Validate PowerPoint file and show its contents"""
    print("=" * 70)
    print("PowerPoint File Validation")
    print("=" * 70)
    print()

    # Check file exists
    if not Path(pptx_path).exists():
        print(f"❌ File not found: {pptx_path}")
        return False

    file_size = Path(pptx_path).stat().st_size
    print(f"✓ File exists: {pptx_path}")
    print(f"  Size: {file_size:,} bytes ({file_size / 1024:.2f} KB)")
    print()

    # Try to open the presentation
    try:
        prs = Presentation(pptx_path)
        print(f"✓ Successfully opened presentation")
        print(f"  Slide dimensions: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")
        print(f"  Total slides: {len(prs.slides)}")
        print()

        # Validate each slide
        print("Slide Contents:")
        print("-" * 70)

        for i, slide in enumerate(prs.slides, 1):
            print(f"\nSlide {i}:")
            print(f"  Shapes: {len(slide.shapes)}")

            # Check for text content
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text[:50])  # First 50 chars

            if text_content:
                print(f"  Text boxes: {len(text_content)}")
                print(f"  First text: \"{text_content[0]}...\"")

            # Check for speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    print(f"  Speaker notes: ✓ ({len(notes_text)} chars)")
                else:
                    print(f"  Speaker notes: (empty)")
            else:
                print(f"  Speaker notes: ✗")

        print()
        print("=" * 70)
        print("✓ Validation SUCCESSFUL - File is valid and readable")
        print("=" * 70)
        return True

    except Exception as e:
        print(f"❌ Error opening presentation: {e}")
        print(f"   Error type: {type(e).__name__}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == '__main__':
    base_dir = Path(__file__).parent
    pptx_path = base_dir / 'output' / 'SoundToAct_TextBased_Presentation.pptx'

    success = validate_pptx(pptx_path)

    if not success:
        print("\n⚠️  The file may be corrupted or invalid.")
        print("   Recommendation: Regenerate the presentation")
