#!/usr/bin/env python3
"""
Detailed PowerPoint inspection
"""

from pathlib import Path
from pptx import Presentation

def inspect_slide_details(pptx_path):
    """Show detailed contents of each slide"""
    prs = Presentation(pptx_path)

    print("=" * 70)
    print("Detailed Slide Inspection")
    print("=" * 70)
    print()

    for i, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*70}")
        print(f"SLIDE {i}")
        print(f"{'='*70}")

        # Get all text content
        print(f"\nText Content ({len(slide.shapes)} shapes):")
        print("-" * 70)

        for j, shape in enumerate(slide.shapes, 1):
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(f"\n[Shape {j}]")
                    print(f"  Type: {shape.shape_type}")
                    print(f"  Text: {text}")

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                print(f"\n{'='*70}")
                print("Speaker Notes:")
                print("-" * 70)
                print(notes_text)

        print()

if __name__ == '__main__':
    base_dir = Path(__file__).parent
    pptx_path = base_dir / 'output' / 'SoundToAct_TextBased_Presentation.pptx'
    inspect_slide_details(pptx_path)
